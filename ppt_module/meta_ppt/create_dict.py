from utils.position import get_position
from pptx.enum.shapes import MSO_SHAPE_TYPE

def create_dict(prs):
    """
    Creates a dictionary representing the structure of a PowerPoint presentation.
    Args:
        prs (Presentation): The PowerPoint presentation object.
    Returns:
        dict: A dictionary containing the structure of the presentation.
        """
    dct={}

    for i, slide in enumerate(prs.slides):
        print(f"visited slide {i}")
        chart ={}
        tab={}
        text={}
        group={}
        img={}

        ids={}
        slide_map={}
        try:
            for shape in slide.shapes:
                #Checking the shape type and assigning it to the corresponding category
                if shape.shape_type ==MSO_SHAPE_TYPE.TABLE:
                    tab[shape.shape_id]= get_position(shape)
                if shape.shape_type==MSO_SHAPE_TYPE.CHART: 
                    chart [shape.shape_id] =get_position(shape)
                if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PLACEHOLDER]:
                    text[shape.shape_id]= get_position(shape) 
                if shape.shape_type== MSO_SHAPE_TYPE.GROUP:
                    group[shape.shape_id] =get_position(shape)
                if shape.shape_type== MSO_SHAPE_TYPE.PICTURE:
                    img[shape.shape_id]= get_position(shape)
                ids[shape.shape_id]=shape
        except Exception as e:
            continue
        # Creating a map of shapes within the slide 
        slide_map["chart"]=chart
        slide_map["table"]=tab
        slide_map["text"]=text
        slide_map["group"]=group
        slide_map["image"]=img
        slide_map["ids"]=ids

        dct[i]=slide_map
    
    return dct