import re
import json
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text(shape):
    """
	Extracts text from a shape object.
	Args:
		shape (Shape): A shape object from a presentation
	Returns:
		str: The extracted text from the shape
	"""
    shape_txt=""
    for paragraph in shape.text_frame.paragraphs:
        shape_txt += paragraph.text.strip()
        if not (shape_txt.endswith(".") or shape_txt.endswith("?") or shape_txt.endswith(":")):
            shape_txt += ". "
    return shape_txt

def extract_table(shape):
    """
    Extracts table data from a shape object.

	Args:
		shape (Shape): A shape object from a presentation
	Returns:
		str: The extracted table data in JSON format
    """
    tab_runs = []
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)

    row_val = []
    spanned_cell_value = ""

    # creating a transposed df
    for c in range(col_count):
        temp = []
        for r in range(col_count):
            cell = tbl.cell(r, c)
            txt = "".join([para.text for para in cell.text_frame.paragraphs])

            if cell.is_merge_origin:
                spanned_cell_value = " ".join([para.text for para in cell.text_frame.paragraphs])
            
            if cell.is_spanned:
                txt = spanned_cell_value

            temp.append(txt)
        tab_runs.append(temp)

    #Converting the 2D List into natural language
    df= pd.DataFrame(tab_runs).transpose().reset_index(drop=True)
    df.columns= df.iloc[0]
    df = df[1:]

    #Return natural language in JSON format
    jsonn = json.dumps(df.to_dict())
    return jsonn

def extract_chart(shape):
    """
    Extracts chart data from a shape object.

    Args:
        shape (Shape): A shape object from a presentation

    Returns:
        str: The extracted chart data in JSON format
    """


    chart=shape.chart
    categories=chart.plots[0].categories
    series=chart.plots[0].series
    data = {}

    # Extract series data and store it in the data dictionary
    for s in series:
        series_name = s.name
        series_values = s.values
        data[series_name] = series_values

    #Create a DataFrame from the extracted data
        try:
            df=pd.DataFrame(data)
            df=df.ffill()
            jsonn=json.dumps(df.to_dict())
            return jsonn
        except Exception as e:
            return json.dumps(data)

def extract_group(shape):
    """
    Extracts information from a group of shapes.
    Args:
        shape (Shape): A shape object representing a group of shapes
    Returns:
        str: The extracted information from the group of shapes
    """

    group_txt=""
    for shp in shape.shapes:
        if shp.has_text_frame:
            group_txt+=shp.text
        if shp.has_table:
            group_txt += extract_table(shp) + "\n"
        if shp.has_chart:
            group_txt += extract_chart(shp) + "\n"
        if shp.shape_type == MSD_SHAPE_TYPE.GROUP:
            group_txt += extract_group(shp)
        else:
            pass
    return group_txt

def extract_data(shp, flag_only_text=False):
    """
    Extracts relevant information from a shape obiect.

    Args:
        shp (Shape): A shape object from a presentation
        flag_only_text (bool, optional): Specifies whether to extract only text. Defaults to False.
    Returns:
        str: The extracted information from the shape object
    """

    if flag_only_text:
    # Extracting only text from the shape if flag_only_text is True
        if shp.has_text_frame:
            try:
                return re.sub("\n", " ", shp.text)
            except:
                return re.sub("\n", " ", shp.text)
    
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                return re.sub("\n", " ", extract_group(shp))
            except:
                return re.sub("\n", " ", extract_group(shp))
        else:
            return ""
    else:
        if shp.has_text_frame:
            try:
                return re.sub("\n", " ", shp.text)
            except:
                return re.sub("\n", " ", shp.text)

        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                return re.sub("\n", " ", extract_group(shp))
            except:
                return re.sub("\n", " ", extract_group(shp))

        if shp.has_table:
            try:
                return re.sub("\n", " ", extract_table(shp))
            except:
                return extract_table(shp)
        if shp.has_chart:
            try:
                return re.sub("\n", " ", extract_chart(shp))
            except:
                return extract_chart(shp)
        else:
            return ""